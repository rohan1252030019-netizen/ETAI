import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors
COLOR_BG = RGBColor(10, 15, 30)         # Very dark blue #0A0F1E
COLOR_PRIMARY = RGBColor(255, 255, 255) # White
COLOR_SECONDARY = RGBColor(170, 190, 210) # Light blue-gray
COLOR_ACCENT = RGBColor(0, 204, 255)    # Cyan
COLOR_CARD = RGBColor(25, 35, 60)       # Glassmorphism base #19233C
COLOR_DARK_CARD = RGBColor(15, 20, 40)
COLOR_SUCCESS = RGBColor(40, 200, 120)  # Green
COLOR_WARNING = RGBColor(255, 170, 0)   # Yellow/Orange

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="PLATFORM VIEW"):
    # Category tracker
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Segoe UI"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY

def add_card(slide, left, top, width, height, bg_color=COLOR_CARD):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = COLOR_ACCENT
    card.line.width = Pt(1)
    return card

def add_bullet_points(slide, left, top, width, height, items, font_size=14, spacing=10):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(spacing)
        p.level = 0

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # ================= SLIDE 1: TITLE SLIDE =================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    
    # Large Logo visual
    logo_bg = s1.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(5.666), Inches(1.5), Inches(2.0), Inches(2.0))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = COLOR_CARD
    logo_bg.line.color.rgb = COLOR_ACCENT
    logo_bg.line.width = Pt(3)
    
    logo_text = s1.shapes.add_textbox(Inches(5.666), Inches(2.1), Inches(2.0), Inches(0.8))
    logo_tf = logo_text.text_frame
    logo_p = logo_tf.paragraphs[0]
    logo_p.text = "IMX"
    logo_p.font.name = "Segoe UI"
    logo_p.font.size = Pt(32)
    logo_p.font.bold = True
    logo_p.font.color.rgb = COLOR_ACCENT
    logo_p.alignment = PP_ALIGN.CENTER

    add_text_box(s1, Inches(1.0), Inches(3.8), Inches(11.333), Inches(0.8), "IMMUNEX", 54, True, COLOR_PRIMARY, PP_ALIGN.CENTER)
    add_text_box(s1, Inches(1.0), Inches(4.7), Inches(11.333), Inches(0.5), "Autonomous AI-Powered Cyber Resilience Platform", 20, False, COLOR_ACCENT, PP_ALIGN.CENTER)
    
    meta_text = "HACKATHON PRESENTATION  •  TEAM ETAI  •  JUNE 2026"
    add_text_box(s1, Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.4), meta_text, 11, True, COLOR_SECONDARY, PP_ALIGN.CENTER)

    # ================= SLIDE 2: PROBLEM STATEMENT =================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "The Critical Vulnerability of Modern Enterprise & CNI", "THE PROBLEM")
    
    problems = [
        ("Reactive Defense", "Security teams wait for telemetry alerts, responding only after a systems compromise has occurred."),
        ("Alert Fatigue", "Thousands of low-context daily alerts cause fatigue, hiding actual multi-stage critical attacks."),
        ("Static Vuln Risk", "CVE prioritization relies on flat CVSS scores rather than dynamic threat actor targeting data."),
        ("Digital Twin Gaps", "No safe playground to model ransomware propagation or SCADA physical manipulation impact."),
        ("Delayed SOAR", "Remediation scripts run blindly without rollback guarantees, risking infrastructure crashes.")
    ]
    
    for idx, (title, desc) in enumerate(problems):
        x = Inches(0.8 + idx * 2.35)
        y = Inches(2.0)
        w = Inches(2.2)
        h = Inches(4.5)
        
        add_card(s2, x, y, w, h, bg_color=COLOR_CARD)
        
        # Red warning accent line at top of card
        line = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 75, 75)
        line.line.fill.background()
        
        # Text
        add_text_box(s2, x + Inches(0.1), y + Inches(0.3), w - Inches(0.2), Inches(0.6), title, 16, True, COLOR_ACCENT)
        add_text_box(s2, x + Inches(0.1), y + Inches(1.0), w - Inches(0.2), Inches(3.2), desc, 12, False, COLOR_SECONDARY)

    # ================= SLIDE 3: VISION =================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "From Detection to Autonomous Prediction", "PLATFORM VISION")
    
    phases = [
        ("DETECT", "Anomaly/Log Ingestion", COLOR_SECONDARY),
        ("CORRELATE", "Multi-stage Alert Fusion", COLOR_SECONDARY),
        ("PREDICT", "GNN-based Attack Paths", COLOR_ACCENT),
        ("SIMULATE", "Digital Twin Safe Replays", COLOR_ACCENT),
        ("MITIGATE", "Rollback-capable SOAR", COLOR_ACCENT),
        ("LEARN", "Vector Feedback Memory", COLOR_SUCCESS)
    ]
    
    for idx, (p_name, p_desc, color) in enumerate(phases):
        x = Inches(0.8 + idx * 1.95)
        y = Inches(3.0)
        w = Inches(1.8)
        h = Inches(2.5)
        
        add_card(s3, x, y, w, h, bg_color=COLOR_CARD)
        
        # Arrows between cards
        if idx < len(phases) - 1:
            arrow = s3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Inches(0.02), y + Inches(1.1), Inches(0.11), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()
            
        add_text_box(s3, x + Inches(0.1), y + Inches(0.3), w - Inches(0.2), Inches(0.6), p_name, 15, True, color, PP_ALIGN.CENTER)
        add_text_box(s3, x + Inches(0.1), y + Inches(1.0), w - Inches(0.2), Inches(1.3), p_desc, 11, False, COLOR_SECONDARY, PP_ALIGN.CENTER)

    # ================= SLIDE 4: ARCHITECTURE OVERVIEW =================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "IMMUNEX Closed-Loop System Architecture", "SYSTEM ARCHITECTURE")
    
    # Draw Layer Columns
    layers = [
        ("1. INGESTION", ["Syslog & Logs", "EDR Telemetry", "Threat Intel Feeds", "ClickHouse Store"], Inches(0.8)),
        ("2. ANALYTICS", ["Threat Fusion Engine", "GNN Path Predictor", "LSTM Sequence Engine", "Neo4j Sync Layer"], Inches(3.8)),
        ("3. DIGITAL TWIN", ["CNI Topologies", "Ransomware Sim", "SCADA Sim", "PageRank Centrals"], Inches(6.8)),
        ("4. REMEDIATION", ["Autonomous Planner", "SOAR Orchestrator", "Transactional Rollback", "FAISS Feedback Cache"], Inches(9.8))
    ]
    
    for name, items, x in layers:
        y = Inches(1.8)
        w = Inches(2.7)
        h = Inches(4.8)
        
        add_card(s4, x, y, w, h)
        add_text_box(s4, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.6), name, 16, True, COLOR_ACCENT)
        
        # Connect items inside card
        item_y = y + Inches(1.0)
        for item in items:
            item_box = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), item_y, w - Inches(0.3), Inches(0.7))
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = COLOR_DARK_CARD
            item_box.line.color.rgb = COLOR_SECONDARY
            item_box.line.width = Pt(0.5)
            
            add_text_box(s4, x + Inches(0.25), item_y + Inches(0.15), w - Inches(0.5), Inches(0.4), item, 12, False, COLOR_PRIMARY, PP_ALIGN.CENTER)
            item_y += Inches(0.9)

    # ================= SLIDE 5: PREDICTIVE ATTACK FORECAST ENGINE =================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Predictive Attack Forecast Engine (PAFE)", "CORE ANALYTICS")
    
    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s5, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Next-Stage Prediction Pipeline", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• GNN Node Classification: Calculates compromised likelihood for downstream assets.",
        "• LSTM Sequence Analyzer: Analyzes telemetry sequences to identify multi-stage attack patterns.",
        "• Markov Chain Transition Matrix: Models threat actor movement probability states.",
        "• Dynamic MITRE TTP Mapping: Forecasts the attacker's next logical tactics and techniques.",
        "• Real-Time Threat Intel Fusion: Intersects asset configuration with external actor profiles."
    ]
    add_bullet_points(s5, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Visual gauge cards on the right
    metrics = [
        ("Next-Asset Compromise Probability", "87.4%", "CRITICAL - DC Primary", RGBColor(255, 75, 75)),
        ("Prediction Confidence (F1 Score)", "0.942", "High Consistency", COLOR_SUCCESS),
        ("Predicted MITRE Technique", "T1021.002", "Remote Service: SMB", COLOR_WARNING)
    ]
    for idx, (m_title, m_val, m_desc, val_color) in enumerate(metrics):
        y = Inches(1.8 + idx * 1.65)
        add_card(s5, Inches(6.7), y, Inches(5.8), Inches(1.50))
        add_text_box(s5, Inches(6.9), y + Inches(0.15), Inches(5.4), Inches(0.4), m_title, 14, True, COLOR_SECONDARY)
        add_text_box(s5, Inches(6.9), y + Inches(0.55), Inches(2.0), Inches(0.6), m_val, 28, True, val_color)
        add_text_box(s5, Inches(9.2), y + Inches(0.65), Inches(3.1), Inches(0.4), m_desc, 13, False, COLOR_PRIMARY)

    # ================= SLIDE 6: ATTACK GRAPH INTELLIGENCE =================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Relational Attack Graph Analytics & Centralities", "GRAPH INTELLIGENCE")
    
    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s6, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Topology Analysis Functions", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Blast Radius Mapping: Quantifies total asset dependencies exposed during compromise.",
        "• Choke Point Detection: Identifies critical network segments that block path progression.",
        "• PageRank Centrality: Computes high-importance asset bridges to structure defense rules.",
        "• Community Clustering: Partitions assets dynamically to isolate lateral zones.",
        "• Neo4j Enterprise Sync: Automatically mirrors local NetworkX twins into production graph stores."
    ]
    add_bullet_points(s6, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Mock Graph visual on the right
    add_card(s6, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8), bg_color=COLOR_DARK_CARD)
    add_text_box(s6, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "Real-Time Attack Path Topology (Visual)", 14, True, COLOR_SECONDARY, PP_ALIGN.CENTER)
    
    # Draw graph nodes & edges using basic shapes
    node_positions = [
        (Inches(7.2), Inches(3.8), "Host-01"),
        (Inches(8.5), Inches(2.8), "EWS-01"),
        (Inches(8.5), Inches(4.8), "OT-Switch"),
        (Inches(9.8), Inches(3.8), "SCADA Master"),
        (Inches(11.2), Inches(3.8), "PLC-Substation", RGBColor(255, 75, 75))
    ]
    # Draw simple connections
    line_coords = [
        (7.2, 3.8, 8.5, 2.8),
        (7.2, 3.8, 8.5, 4.8),
        (8.5, 2.8, 9.8, 3.8),
        (8.5, 4.8, 9.8, 3.8),
        (9.8, 3.8, 11.2, 3.8)
    ]
    for x1, y1, x2, y2 in line_coords:
        connector = s6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x1+0.2), Inches(y1+0.2), Inches(x2-x1-0.2), Inches(y2-y1 if y2 != y1 else 0.05))
        connector.fill.solid()
        connector.fill.fore_color.rgb = COLOR_ACCENT
        connector.line.fill.background()
        
    for idx, pos in enumerate(node_positions):
        cx, cy, label = pos[0], pos[1], pos[2]
        color = pos[3] if len(pos) > 3 else COLOR_CARD
        circle = s6.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = COLOR_ACCENT
        circle.line.width = Pt(1.5)
        add_text_box(s6, cx - Inches(0.3), cy + Inches(0.65), Inches(1.2), Inches(0.3), label, 9, True, COLOR_PRIMARY, PP_ALIGN.CENTER)

    # ================= SLIDE 7: DIGITAL TWIN SIMULATOR =================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Multi-Sector Digital Twin Simulator", "DIGITAL TWIN")
    
    sectors = [
        ("Energy Grid", "Simulates Modbus/DNP3 commands, tripping circuit breakers and SCADA systems.", "EG-SCADA-SRV-01"),
        ("Government", "Replays lateral pivots through VPN gateways and domain controllers.", "GOV-DC-01"),
        ("Healthcare", "Tracks ransomware propagation exposing medical IoT pumps and EHR storage.", "HC-EHR-DB"),
        ("Education", "Models campus cores, LMS servers, and research GPU clusters.", "EDU-LMS-SRV"),
        ("Telecom", "Forecasts routing loops and cellular core infrastructure signaling failures.", "TEL-CORE-RT"),
        ("Finance", "Replays attacks targeting high-frequency transaction databases.", "FIN-DB-TX")
    ]
    
    for idx, (name, desc, entry) in enumerate(sectors):
        col = idx % 3
        row = idx // 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.8 + row * 2.5)
        w = Inches(3.6)
        h = Inches(2.2)
        
        add_card(s7, x, y, w, h)
        add_text_box(s7, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.4), name, 16, True, COLOR_ACCENT)
        add_text_box(s7, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(1.1), desc, 11, False, COLOR_SECONDARY)
        add_text_box(s7, x + Inches(0.15), y + Inches(1.70), w - Inches(0.3), Inches(0.4), f"Default Patient-0: {entry}", 10, True, COLOR_PRIMARY)

    # ================= SLIDE 8: THREAT ACTOR INTELLIGENCE =================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Threat Actor Intelligence & Attribution", "THREAT INTELLIGENCE")
    
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s8, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Attribution & Profiling Capabilities", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Campaign Correlation: Group isolated indicators into active campaign models.",
        "• Infrastructure Mapping: Resolves IP address reuse and DNS dynamic hosts.",
        "• TTP Clustering: Clusters attack patterns against MITRE ATT&CK groups.",
        "• Actor Profiling: Dynamically fetches target sector intent models.",
        "• Live Threat Feed Sync: Keeps local vulnerability scoring aligned to active actors."
    ]
    add_bullet_points(s8, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Table of current active threats
    add_card(s8, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8), bg_color=COLOR_DARK_CARD)
    add_text_box(s8, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "Tracked CNI Threat Groups", 14, True, COLOR_PRIMARY, PP_ALIGN.CENTER)
    
    # Draw simple table structure
    tx_table = s8.shapes.add_textbox(Inches(6.9), Inches(2.6), Inches(5.4), Inches(3.5))
    tf = tx_table.text_frame
    tf.word_wrap = True
    headers = "Group        | Targets       | Primary Vector     | Status"
    rows = [
        "APT29        | Gov, Energy   | OAuth Abuse        | Active",
        "Volt Typhoon | CNI, Power    | VPN Exploit        | Critical",
        "LockBit 3.0  | Healthcare    | RDP / Phishing     | High",
        "Lazarus Group| Finance       | Supply Chain       | Active"
    ]
    
    p = tf.paragraphs[0]
    p.text = headers
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for row in rows:
        p_row = tf.add_paragraph()
        p_row.text = row
        p_row.font.name = "Consolas"
        p_row.font.size = Pt(11)
        p_row.font.color.rgb = COLOR_PRIMARY
        p_row.space_before = Pt(15)

    # ================= SLIDE 9: AUTONOMOUS MITIGATION PLANNER =================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Autonomous Mitigation Planner (AMP)", "REMEDIATION DECISIONS")
    
    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s9, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Decision Matrix Optimization", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Dynamic Trade-offs: Computes metrics balances between protection levels and costs.",
        "• Resource Constraint Rules: Limits plan scopes depending on engineers bandwidth.",
        "• Asset Dependency Isolation: Prioritizes host network segments to prevent pivot cascades.",
        "• Automatic Feedback: Retrains plans based on playbook performance outcomes.",
        "• Explainable Rationale: Outputs detailed mathematical proofs for each choice."
    ]
    add_bullet_points(s9, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Recommendation visual card
    add_card(s9, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8))
    add_text_box(s9, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "AMP Optimal Action Plan (Forensics)", 14, True, COLOR_SECONDARY, PP_ALIGN.CENTER)
    
    rec_items = [
        ("1. Isolate Host 'OT-GW-01'", "Risk Reduction: -64% | Cost: Low | Speed: <1s", COLOR_SUCCESS),
        ("2. Block Protocol 'Modbus/TCP'", "Risk Reduction: -22% | Cost: Med | Speed: <1s", COLOR_WARNING),
        ("3. Patch DC Controller CVE-2024", "Risk Reduction: -94% | Cost: High | Speed: 2h", COLOR_PRIMARY),
    ]
    for idx, (title, desc, color) in enumerate(rec_items):
        item_y = Inches(2.6 + idx * 1.3)
        add_card(s9, Inches(6.9), item_y, Inches(5.4), Inches(1.0), bg_color=COLOR_DARK_CARD)
        add_text_box(s9, Inches(7.1), item_y + Inches(0.15), Inches(5.0), Inches(0.3), title, 14, True, color)
        add_text_box(s9, Inches(7.1), item_y + Inches(0.50), Inches(5.0), Inches(0.3), desc, 11, False, COLOR_SECONDARY)

    # ================= SLIDE 10: NATIONAL CYBER RESILIENCE INDEX =================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "National Cyber Resilience Index (NCRI)", "EXECUTIVE DASHBOARD")
    
    add_card(s10, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s10, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Measuring National Readiness", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Real-Time Resilience KPI: Scores infrastructure readiness on 0-100 scales.",
        "• Dynamic Weight Modeling: Multiplies sector impacts based on dependency weights.",
        "• Strategic Trend Forecasting: Forecasts scoring changes using historical log trends.",
        "• Policy Alignment Check: Cross-checks national security requirements.",
        "• Executive Reporting: Provides CISO dashboards with single-button exports."
    ]
    add_bullet_points(s10, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Large Score Display on right
    add_card(s10, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8))
    add_text_box(s10, Inches(6.9), Inches(2.2), Inches(5.4), Inches(0.4), "Current National Resilience Level", 16, True, COLOR_SECONDARY, PP_ALIGN.CENTER)
    add_text_box(s10, Inches(6.9), Inches(2.8), Inches(5.4), Inches(1.5), "82.5", 96, True, COLOR_ACCENT, PP_ALIGN.CENTER)
    add_text_box(s10, Inches(6.9), Inches(4.3), Inches(5.4), Inches(0.4), "STATUS: RESILIENT (Up 2.4% over 30 days)", 14, True, COLOR_SUCCESS, PP_ALIGN.CENTER)
    
    add_card(s10, Inches(7.2), Inches(5.0), Inches(4.8), Inches(1.1), bg_color=COLOR_DARK_CARD)
    add_text_box(s10, Inches(7.4), Inches(5.15), Inches(4.4), Inches(0.8), "Recommendation: Segment OT-Switch zone from IT-Domain controller networks to lift NCRI to 86.8.", 12, False, COLOR_SECONDARY, PP_ALIGN.CENTER)

    # ================= SLIDE 11: EXPLAINABLE AI =================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "Explainable AI (XAI) Risk Explainer", "TRANSPARENCY & AUDITING")
    
    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s11, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Demystifying AI Actions", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Evidence Traceability: Every PAFE prediction attaches exact logs and feed sources.",
        "• Risk Decomposition: Breaks dynamic scores down to CVSS, KEV, and actor weights.",
        "• MITRE Alignment Details: Maps suggested mitigations directly to target TTP identifiers.",
        "• Visual Graph Rationale: highlights paths traversed by simulation algorithms.",
        "• No Black Box Models: Guarantees auditor clarity for all autonomous defensive playbooks."
    ]
    add_bullet_points(s11, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Audit log entry visual on right
    add_card(s11, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8), bg_color=COLOR_DARK_CARD)
    add_text_box(s11, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "XAI Explanation Log (Example)", 14, True, COLOR_ACCENT, PP_ALIGN.CENTER)
    
    xai_log = s11.shapes.add_textbox(Inches(6.9), Inches(2.5), Inches(5.4), Inches(3.7))
    tf_xai = xai_log.text_frame
    tf_xai.word_wrap = True
    
    p = tf_xai.paragraphs[0]
    p.text = "EVENT: Target DC-01 compromise probability forecasted at 87%."
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    logs = [
        "EVIDENCE 1: Active Volt Typhoon actor observed targets.",
        "EVIDENCE 2: CVE-2024-3094 vulnerability active on segment.",
        "EVIDENCE 3: Direct network edges found between Client-01 and DC-01.",
        "EXPLANATION: LIME coefficients: CVSS (0.35) + KEV (0.45) + Network-Hops (0.20) = 1.0. Mitigation suggested: isolate Client-01 immediately."
    ]
    for line in logs:
        p_line = tf_xai.add_paragraph()
        p_line.text = line
        p_line.font.name = "Consolas"
        p_line.font.size = Pt(10)
        p_line.font.color.rgb = COLOR_SECONDARY
        p_line.space_before = Pt(12)

    # ================= SLIDE 12: CYBER LEARNING MEMORY =================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_header(s12, "Cyber Learning Memory Loop", "FEEDBACK OPTIMIZATION")
    
    add_card(s12, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s12, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Vector-Powered Knowledge Retention", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Incident Learning Database: Indexes mitigation outcomes in vector stores (FAISS).",
        "• Automated Retrieval: Queries similar past attacks to suggest tested playbooks.",
        "• Continuous Optimization: Tracks mitigation metrics to optimize weights.",
        "• Edge-Store Resilience: Restores database structures from localized binary caches.",
        "• Shared Defense Models: Enables secure cross-sector indicators synchronizations."
    ]
    add_bullet_points(s12, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Process loop on the right
    add_card(s12, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8))
    add_text_box(s12, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "Closed-Loop Learning Cycle", 14, True, COLOR_SECONDARY, PP_ALIGN.CENTER)
    
    steps = [
        ("1. Ingestion & Anomaly Score", "Telemetry evaluated, anomaly triggers.", Inches(2.5)),
        ("2. Vector Similarity Search", "FAISS query finds matching past threat states.", Inches(3.7)),
        ("3. Adaptive Plan RETUNE", "Mitigation weights updated by feedback memory.", Inches(4.9))
    ]
    for s_title, s_desc, s_y in steps:
        add_card(s12, Inches(6.9), s_y, Inches(5.4), Inches(1.0), bg_color=COLOR_DARK_CARD)
        add_text_box(s12, Inches(7.1), s_y + Inches(0.15), Inches(5.0), Inches(0.3), s_title, 14, True, COLOR_ACCENT)
        add_text_box(s12, Inches(7.1), s_y + Inches(0.50), Inches(5.0), Inches(0.3), s_desc, 11, False, COLOR_PRIMARY)

    # ================= SLIDE 13: SOC COPILOT =================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13)
    add_header(s13, "IMMUNEX SOC Copilot Chat Interface", "HUMAN-IN-THE-LOOP")
    
    add_card(s13, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    add_text_box(s13, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5), "Natural Language Security Operations", 18, True, COLOR_ACCENT)
    
    bullets = [
        "• Threat Analysis Prompts: Queries platform state using conversational English.",
        "• Automated Rule Synthesis: Automatically compiles Sigma or YARA rules from chat.",
        "• Playbook Control Queries: Inspects SOAR states and launches manual rollbacks.",
        "• MITRE Technique Summaries: Explains complex multi-stage attack steps.",
        "• Seamless Integration: Links into enterprise messaging hubs (Slack webhook)."
    ]
    add_bullet_points(s13, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.8), bullets, 13, 14)
    
    # Mock Chat Window
    add_card(s13, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8), bg_color=COLOR_DARK_CARD)
    add_text_box(s13, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4), "SOC Copilot Console (Mockup)", 14, True, COLOR_PRIMARY, PP_ALIGN.CENTER)
    
    # User message
    add_card(s13, Inches(8.5), Inches(2.5), Inches(3.8), Inches(0.8), bg_color=COLOR_CARD)
    add_text_box(s13, Inches(8.6), Inches(2.55), Inches(3.6), Inches(0.7), "User: Is there any Volt Typhoon lateral movement threat to EG-PLC-01?", 10, False, COLOR_SECONDARY)
    
    # Copilot Response
    add_card(s13, Inches(6.9), Inches(3.5), Inches(4.8), Inches(1.6))
    copilot_msg = (
        "Copilot: WARNING. GNN attack path model predicts 87.4% chance of "
        "lateral pivot through EG-DMZ-JUMP to EG-PLC-01 using Kerberos delegation.\n\n"
        "Mitigation suggested: run playbooks/scada_tamper.yml immediately to isolate EG-DMZ-JUMP."
    )
    add_text_box(s13, Inches(7.0), Inches(3.55), Inches(4.6), Inches(1.5), copilot_msg, 10, False, COLOR_PRIMARY)
    
    # Input area
    add_card(s13, Inches(6.9), Inches(5.4), Inches(5.4), Inches(0.8), bg_color=COLOR_DARK_CARD)
    add_text_box(s13, Inches(7.0), Inches(5.6), Inches(5.2), Inches(0.4), "Type a message or command (e.g., /mitigate eg-dmz-jump)...", 11, False, COLOR_SECONDARY)

    # ================= SLIDE 14: ENTERPRISE SECURITY CONTROLS =================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14)
    add_header(s14, "Enterprise Security & Platform Controls", "PLATFORM SECURITY")
    
    controls = [
        ("Role-Based Access Control", "Enforces fine-grained JWT/OAuth scope checks across all API requests and command interfaces."),
        ("Cryptographic Audit Logs", "Maintains local, append-only verification databases ensuring unalterable recording of all playbook actions."),
        ("API Versioning & Limits", "Implements strict rate limits and version management to shield analytical models from DoS events."),
        ("Observability & Telemetry", "Outputs metrics through Prometheus, Grafana, and Loguru pipelines, keeping track of ingestion health.")
    ]
    
    for idx, (title, desc) in enumerate(controls):
        col = idx % 2
        row = idx // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.5)
        w = Inches(5.6)
        h = Inches(2.2)
        
        add_card(s14, x, y, w, h)
        add_text_box(s14, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.4), title, 16, True, COLOR_ACCENT)
        add_text_box(s14, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(1.3), desc, 12, False, COLOR_SECONDARY)

    # ================= SLIDE 15: PERFORMANCE METRICS =================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_background(s15)
    add_header(s15, "IMMUNEX Platform Performance Metrics", "VALIDATION STATS")
    
    kpis = [
        ("Verification Test Suite", "587 / 587", "Tests Passed Successfully (100%)", COLOR_SUCCESS),
        ("Detection Anomaly Speed", "< 12 ms", "Average telemetry sequence analysis latency", COLOR_ACCENT),
        ("Path Prediction Accuracy", "94.2%", "F1-score across CNI topology datasets", COLOR_ACCENT),
        ("Mean Time to Respond", "< 1 s", "SOAR playbook trigger-to-mitigation time", COLOR_SUCCESS),
        ("Blast Radius Reduction", "-72.4%", "Average damage reduction under AMP controls", COLOR_SUCCESS),
        ("Active CVE Risk Mitigation", "-94.8%", "Vulnerability risk drops post optimal patching", COLOR_SUCCESS)
    ]
    
    for idx, (title, val, desc, val_color) in enumerate(kpis):
        col = idx % 3
        row = idx // 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.8 + row * 2.5)
        w = Inches(3.6)
        h = Inches(2.2)
        
        add_card(s15, x, y, w, h)
        add_text_box(s15, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.40), title, 14, True, COLOR_SECONDARY)
        add_text_box(s15, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.7), val, 32, True, val_color)
        add_text_box(s15, x + Inches(0.15), y + Inches(1.35), w - Inches(0.3), Inches(0.7), desc, 11, False, COLOR_PRIMARY)

    # ================= SLIDE 16: COMPETITIVE COMPARISON =================
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_background(s16)
    add_header(s16, "Competitive Feature Matrix", "MARKET ANALYSIS")
    
    # Table layout
    table_x = Inches(0.8)
    table_y = Inches(1.8)
    table_w = Inches(11.733)
    table_h = Inches(4.8)
    
    add_card(s16, table_x, table_y, table_w, table_h)
    
    # Add comparison content as structured text
    txBox = s16.shapes.add_textbox(table_x + Inches(0.2), table_y + Inches(0.2), table_w - Inches(0.4), table_h - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    headers = "Feature Capability            | SIEM/SOAR | XDR       | ASM       | IMMUNEX (AI Platform)"
    rows = [
        "Dynamic Attack Graph Prediction  | No        | No        | Static    | Yes (Real-Time GNN Path)",
        "Pluggable Digital Twin Simulator | No        | No        | No        | Yes (Sectors EG, HC, GOV)",
        "Vulnerability Actor Context      | Static    | Static    | Partial   | Yes (KEV + Actor Target Sync)",
        "Playbook Action Rollback Safety  | No        | No        | No        | Yes (Transactional Undo)",
        "Cyber Learning Vector Memory     | No        | No        | No        | Yes (FAISS Index Loops)",
        "Executive Resilience Scoring    | No        | No        | No        | Yes (Live NCRI Score)"
    ]
    
    p = tf.paragraphs[0]
    p.text = headers
    p.font.name = "Consolas"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for row in rows:
        p_row = tf.add_paragraph()
        p_row.text = row
        p_row.font.name = "Consolas"
        p_row.font.size = Pt(11)
        p_row.font.color.rgb = COLOR_PRIMARY
        p_row.space_before = Pt(20)

    # ================= SLIDE 17: DEMO WORKFLOW =================
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_background(s17)
    add_header(s17, "Closed-Loop Incident Scenario Demo", "PLATFORM DEMO")
    
    steps = [
        ("1. Ingestion", "Volt Typhoon credential abuse observed at firewall.", Inches(0.8)),
        ("2. Fusion", "Threat feed correlates IP with known active command server.", Inches(2.8)),
        ("3. Prediction", "PAFE engine forecasts 87% compromise path to classified DB.", Inches(4.8)),
        ("4. Simulation", "Digital Twin simulates domain controller pivot scenarios.", Inches(6.8)),
        ("5. Mitigation", "SOAR executes sandbox isolation with rollback guards.", Inches(8.8)),
        ("6. Learning", "Mitigation vector indexed in FAISS memory store for retuning.", Inches(10.8))
    ]
    
    for title, desc, x in steps:
        y = Inches(2.2)
        w = Inches(1.8)
        h = Inches(4.0)
        
        add_card(s17, x, y, w, h)
        add_text_box(s17, x + Inches(0.1), y + Inches(0.2), w - Inches(0.2), Inches(0.6), title, 15, True, COLOR_ACCENT)
        add_text_box(s17, x + Inches(0.1), y + Inches(0.9), w - Inches(0.2), Inches(2.9), desc, 11, False, COLOR_SECONDARY)

    # ================= SLIDE 18: BUSINESS IMPACT =================
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_background(s18)
    add_header(s18, "Strategic Impact & Business Outcomes", "BUSINESS CASE")
    
    impacts = [
        ("Minimize Downtime", "Reduces potential outage impacts by 72% across critical services (electricity, medical, telecom) under automated, fast remediations."),
        ("Optimize Budgets", "Contextual risk prioritizations help security leads focus budgets on critical vulnerabilities, lowering operational costs."),
        ("Increase Compliance", "Provides unalterable, automated log trails verifying operations conform with ISO/IEC 27001 and national rules."),
        ("Strengthen Insurance", "Provides underwriter consoles with real-time NCRI scores, qualifying enterprises for cheaper cybersecurity insurance plans.")
    ]
    
    for idx, (title, desc) in enumerate(impacts):
        col = idx % 2
        row = idx // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.5)
        w = Inches(5.6)
        h = Inches(2.2)
        
        add_card(s18, x, y, w, h)
        add_text_box(s18, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.4), title, 16, True, COLOR_ACCENT)
        add_text_box(s18, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(1.3), desc, 12, False, COLOR_SECONDARY)

    # ================= SLIDE 19: FUTURE ROADMAP =================
    s19 = prs.slides.add_slide(blank_layout)
    set_slide_background(s19)
    add_header(s19, "IMMUNEX Strategic Roadmap", "THE FUTURE")
    
    phases = [
        ("Phase 1: Q3-Q4 2026", "Federated Defense Engine", "Synchronizes threat actor memory models across air-gapped grids without leaking local network details."),
        ("Phase 2: Q1-Q2 2027", "Hardware-in-the-Loop Sim", "Deploys direct simulated network signals to physical industrial controllers (PLCs, RTUs)."),
        ("Phase 3: Q3-Q4 2027", "Generative Playbooks", "Generates safe, rollback-validated remediations dynamically using LLMs tuned on security records.")
    ]
    
    for idx, (title, name, desc) in enumerate(phases):
        y = Inches(1.8 + idx * 1.65)
        add_card(s19, Inches(0.8), y, Inches(11.733), Inches(1.4))
        add_text_box(s19, Inches(1.0), y + Inches(0.15), Inches(3.5), Inches(0.3), title, 14, True, COLOR_ACCENT)
        add_text_box(s19, Inches(1.0), y + Inches(0.50), Inches(3.5), Inches(0.7), name, 16, True, COLOR_PRIMARY)
        add_text_box(s19, Inches(4.8), y + Inches(0.25), Inches(7.5), Inches(0.9), desc, 12, False, COLOR_SECONDARY)

    # ================= SLIDE 20: CLOSING SLIDE =================
    s20 = prs.slides.add_slide(blank_layout)
    set_slide_background(s20)
    
    add_text_box(s20, Inches(1.0), Inches(2.2), Inches(11.333), Inches(0.8), "The Future of Cyber Resilience Starts Here", 36, True, COLOR_ACCENT, PP_ALIGN.CENTER)
    add_text_box(s20, Inches(1.0), Inches(3.2), Inches(11.333), Inches(1.0), "IMMUNEX", 64, True, COLOR_PRIMARY, PP_ALIGN.CENTER)
    add_text_box(s20, Inches(1.0), Inches(4.5), Inches(11.333), Inches(0.4), "Autonomous AI-Powered Security for Modern Critical Infrastructure", 16, False, COLOR_SECONDARY, PP_ALIGN.CENTER)
    
    add_text_box(s20, Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.5), "THANK YOU", 20, True, COLOR_ACCENT, PP_ALIGN.CENTER)

    prs.save("IMMUNEX_Hackathon_Final.pptx")
    print("Presentation created successfully!")

if __name__ == '__main__':
    create_presentation()
